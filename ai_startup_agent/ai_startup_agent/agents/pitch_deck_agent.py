"""agents/pitch_deck_agent.py — Pitch Deck Generation Agent

Turns the consolidated analysis into a downloadable .pptx investor deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8A, 0x8F, 0x98)


class PitchDeckAgent:
    name = "Pitch Deck Agent"

    def run(self, idea_title: str, all_data: dict) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        self._title_slide(prs, blank, idea_title, all_data)
        self._market_slide(prs, blank, all_data.get("market", {}))
        self._competitor_slide(prs, blank, all_data.get("competitor", {}))
        self._technical_slide(prs, blank, all_data.get("technical", {}))
        self._business_model_slide(prs, blank, all_data.get("business_model", {}))
        self._financial_slide(prs, blank, all_data.get("financial", {}))
        self._validation_slide(prs, blank, all_data.get("validation", {}))

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ---- slide builders ----------------------------------------------

    def _bg(self, slide, color=NAVY):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _add_title(self, slide, text, top=0.5):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(1))
        tf = box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = WHITE
        return box

    def _add_bullets(self, slide, items, top=1.6, left=0.6, width=12, size=18):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = WHITE
            p.space_after = Pt(10)

    def _title_slide(self, prs, layout, idea_title, all_data):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2))
        tf = box.text_frame
        tf.text = idea_title
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        verdict = all_data.get("validation", {}).get("verdict", "")
        score = all_data.get("validation", {}).get("overall_score", "")
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(1))
        tf2 = sub.text_frame
        tf2.text = f"AI Viability Verdict: {verdict}  |  Score: {score}/10"
        tf2.paragraphs[0].font.size = Pt(20)
        tf2.paragraphs[0].font.color.rgb = ACCENT

    def _market_slide(self, prs, layout, market):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        self._add_title(slide, "Market Opportunity")
        bullets = [
            f"TAM: ${market.get('tam_usd_billion', '?')}B | SAM: ${market.get('sam_usd_billion', '?')}B | SOM: ${market.get('som_usd_billion', '?')}B",
            f"CAGR: {market.get('cagr_percent', '?')}%",
        ] + market.get("key_trends", [])
        self._add_bullets(slide, bullets)

    def _competitor_slide(self, prs, layout, competitor):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        self._add_title(slide, "Competitive Landscape")
        bullets = [f"{c.get('name')} ({c.get('type')}) — positioning {c.get('positioning_score')}/10"
                   for c in competitor.get("competitors", [])]
        bullets.append(f"Market gap: {competitor.get('market_gap', '')}")
        self._add_bullets(slide, bullets, size=16)

    def _technical_slide(self, prs, layout, technical):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        self._add_title(slide, "Technical Feasibility")
        bullets = [
            f"Feasibility score: {technical.get('feasibility_score', '?')}/10",
            f"MVP timeline: {technical.get('mvp_timeline_weeks', '?')} weeks with {technical.get('team_size_needed', '?')} engineers",
            "Stack: " + ", ".join(technical.get("recommended_stack", [])),
        ] + technical.get("key_technical_risks", [])
        self._add_bullets(slide, bullets)

    def _business_model_slide(self, prs, layout, bm):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        self._add_title(slide, "Business Model & GTM")
        bullets = [f"Model: {bm.get('business_model', '')}"]
        bullets += [r.get("name", "") for r in bm.get("revenue_streams", [])]
        bullets += bm.get("gtm_strategy", {}).get("primary_channels", [])
        self._add_bullets(slide, bullets)

    def _financial_slide(self, prs, layout, fin):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        self._add_title(slide, "Financial Projections")
        bullets = []
        for y in fin.get("yearly_projections", []):
            bullets.append(
                f"Year {y.get('year')}: ${y.get('revenue_usd'):,} revenue, {y.get('customers')} customers"
                if isinstance(y.get("revenue_usd"), (int, float)) else str(y)
            )
        bullets.append(f"Funding needed: ${fin.get('initial_funding_needed_usd', '?'):,}"
                        if isinstance(fin.get("initial_funding_needed_usd"), (int, float)) else "")
        self._add_bullets(slide, [b for b in bullets if b])

    def _validation_slide(self, prs, layout, val):
        slide = prs.slides.add_slide(layout)
        self._bg(slide)
        self._add_title(slide, "Final Verdict")
        bullets = [f"Overall score: {val.get('overall_score', '?')}/10 — {val.get('verdict', '')}"]
        bullets += ["Strength: " + s for s in val.get("top_strengths", [])]
        bullets += ["Risk: " + r for r in val.get("top_risks", [])]
        self._add_bullets(slide, bullets)
